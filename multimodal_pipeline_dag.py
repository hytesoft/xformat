# 保证本地模块可被导入
import sys
import yaml
import mimetypes
import os

# 加载全局配置
with open('config/config.yaml', 'r') as f:
    _config = yaml.safe_load(f)
PUBLIC_DIR = _config.get('public_dir', './public')

# 所有节点函数定义（scan_file_task、classify_file_task、segment_text_task、...、video_multimodal_understanding_task）全部保留

def scan_file_task(input_file, output_json=None):
    print(f"[SCAN] 开始扫描文件/目录: {input_file}")
    import os
    from io_utils.file_scanner import scan_files
    if not output_json:
        output_json = _config.get("output_json", os.path.join(PUBLIC_DIR, "file_list.json"))
    scan_files(input_file, output_json)
    print(f"[SCAN] 扫描完成，结果写入: {output_json}")

def classify_file_task(input_file):
    """
    判断文件类型，返回分段节点名。无Airflow依赖，参数直接传递。
    """
    import mimetypes
    filetype, _ = mimetypes.guess_type(input_file)
    if not filetype:
        # 针对常见 office 文件补充判断
        if input_file.lower().endswith('.docx'):
            return 'segment_word'
        if input_file.lower().endswith('.pptx'):
            return 'segment_ppt'
        if input_file.lower().endswith('.xlsx'):
            return 'segment_excel'
        raise ValueError(f"无法识别文件类型: {input_file}")
    # 文本类类型判断
    if filetype.startswith('text/plain') or filetype.startswith('text/markdown') or filetype.startswith('text/html'):
        return 'segment_text'
    if filetype.startswith('application/pdf'):
        return 'segment_pdf'
    if filetype in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
        return 'segment_word'
    if filetype in ['application/vnd.openxmlformats-officedocument.presentationml.presentation']:
        return 'segment_ppt'
    if filetype in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
        return 'segment_excel'
    elif filetype.startswith('audio/'):
        return 'segment_audio'
    elif filetype.startswith('image/'):
        return 'segment_image'
    elif filetype.startswith('video/'):
        return 'segment_video'
    else:
        return 'segment_other'

def segment_text_task(input_file, output_path):
    print(f"[TEXT] 开始分段文本文件: {input_file}")
    import os, json, re
    from bs4 import BeautifulSoup
    import mimetypes
    filetype, _ = mimetypes.guess_type(input_file)
    segments = []
    with open(input_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    for idx, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        # html: 去标签、去样式、去脚本
        if filetype and filetype.startswith('text/html'):
            soup = BeautifulSoup(line, 'html.parser')
            text = soup.get_text(separator=' ', strip=True)
        # markdown: 去除常见语法符号
        elif filetype and filetype.startswith('text/markdown'):
            text = re.sub(r'[#*`>\-\[\]()!_~]', '', line)
            text = re.sub(r'\!\[.*?\]\(.*?\)', '', text)  # 去图片
            text = re.sub(r'\[.*?\]\(.*?\)', '', text)    # 去链接
            text = text.strip()
        # 纯文本：去除特殊控制字符
        else:
            text = re.sub(r'[\x00-\x1f\x7f]', '', line)
        if text:
            segments.append({
                "index": idx,
                "text": text,
                "audio": None,
                "image": []
            })
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[TEXT] 分段完成，共{len(segments)}段，输出: {output_path}")

def segment_pdf_task(input_file, output_path):
    print(f"[PDF] 开始分段PDF文件: {input_file}")
    import os, json
    import fitz  # PyMuPDF
    doc = fitz.open(input_file)
    segments = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text().strip()
        images = []
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = os.path.join(PUBLIC_DIR, f"pdf_page{page_num+1}_img{img_index+1}.{image_ext}")
            with open(image_filename, "wb") as img_f:
                img_f.write(image_bytes)
            images.append(image_filename)
        # --- OCR兜底逻辑 ---
        if not text or len(text) < 10:
            try:
                import pytesseract
                from PIL import Image
                pix = page.get_pixmap()
                ocr_img_path = os.path.join(PUBLIC_DIR, f"pdf_page{page_num+1}_ocr.png")
                pix.save(ocr_img_path)
                # 兼容中英文，优先用chi_sim，失败则退回eng
                try:
                    ocr_text = pytesseract.image_to_string(Image.open(ocr_img_path), lang="chi_sim+eng")
                except Exception:
                    ocr_text = pytesseract.image_to_string(Image.open(ocr_img_path), lang="eng")
                ocr_text = ocr_text.strip()
                if ocr_text:
                    if text:
                        text = text + "\n[OCR补全]:\n" + ocr_text
                    else:
                        text = ocr_text
                images.append(ocr_img_path)
            except Exception as e:
                print(f"[OCR兜底失败] page {page_num+1}: {e}")
        segments.append({
            "index": page_num,
            "text": text,
            "audio": None,
            "image": images
        })
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[PDF] 分段完成，共{len(segments)}页，输出: {output_path}")

def segment_word_task(input_file, output_path):
    print(f"[WORD] 开始分段Word文件: {input_file}")
    import os, json
    from docx import Document
    doc = Document(input_file)
    segments = []
    img_idx = 0
    for idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        images = []
        # 提取图片（docx图片一般在 runs 里）
        for run in para.runs:
            if 'graphic' in run._element.xml:
                # 这里只能做简单标记，复杂图片提取需用更底层包
                images.append(f"word_img_{img_idx+1}")
                img_idx += 1
        if text or images:
            segments.append({
                "index": idx,
                "text": text,
                "audio": None,
                "image": images
            })
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[WORD] 分段完成，共{len(segments)}段，输出: {output_path}")

def segment_ppt_task(input_file, output_path):
    print(f"[PPT] 开始分段PPT文件: {input_file}")
    import os, json
    from pptx import Presentation
    prs = Presentation(input_file)
    segments = []
    img_idx = 0
    for idx, slide in enumerate(prs.slides):
        texts = []
        images = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
            if shape.shape_type == 13:  # PICTURE
                images.append(f"ppt_img_{img_idx+1}")
                img_idx += 1
        text = "\n".join([t for t in texts if t])
        segments.append({
            "index": idx,
            "text": text,
            "audio": None,
            "image": images
        })
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[PPT] 分段完成，共{len(segments)}页，输出: {output_path}")

def segment_excel_task(input_file, output_path):
    print(f"[EXCEL] 开始分段Excel文件: {input_file}")
    import os, json
    import pandas as pd
    xls = pd.ExcelFile(input_file)
    segments = []
    idx = 0
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        for row in df.itertuples(index=False):
            text = ' '.join([str(cell) for cell in row if pd.notnull(cell)])
            if text:
                segments.append({
                    "index": idx,
                    "text": text,
                    "audio": None,
                    "image": []
                })
                idx += 1
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[EXCEL] 分段完成，共{len(segments)}行，输出: {output_path}")

def segment_audio_task(input_file, output_path):
    print(f"[AUDIO] 开始分段音频文件: {input_file}")
    # 本节点只负责音频分段，不做音频转文字（ASR），所有音频转文字由 asr_whisper_for_segments 统一处理。
    import os, json, math
    import yaml
    from pydub import AudioSegment
    import mimetypes
    filetype, _ = mimetypes.guess_type(input_file)
    if not (filetype and filetype.startswith('audio/')):
        raise ValueError(f"segment_audio_task 只处理音频文件，收到非音频文件: {input_file}, type: {filetype}")
    # 读取配置
    config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    segment_seconds = config.get('audio_segment_seconds', 30)
    audio = AudioSegment.from_file(input_file)
    duration_sec = len(audio) / 1000
    num_segments = math.ceil(duration_sec / segment_seconds)
    segments = []
    for idx in range(num_segments):
        start_ms = idx * segment_seconds * 1000
        end_ms = min((idx + 1) * segment_seconds * 1000, len(audio))
        segment_audio = audio[start_ms:end_ms]
        seg_path = os.path.join(PUBLIC_DIR, f'audio_seg_{idx+1}.wav')
        segment_audio.export(seg_path, format='wav')
        segments.append({
            "index": idx,
            "text": None,
            "audio": seg_path,
            "image": []
        })
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[AUDIO] 分段完成，共{len(segments)}段，输出: {output_path}")

def segment_video_task(input_file, output_path):
    print(f"[VIDEO] 开始分段视频文件: {input_file}")
    """
    本节点只负责视频分段（定时抽帧+音频切片），不做音频转文字（ASR），所有音频转文字由 asr_whisper_for_segments 统一处理。
    每隔 config['video_force_segment_sec'] 秒强制切一帧图片，并提取对应音频片段，兜底输出首帧和首段音频，输出结构统一包含 index/text/audio/image 四字段。
    """
    import os, json, yaml, cv2
    import math
    import subprocess
    config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    force_segment_sec = config.get('video_force_segment_sec', 15)
    cap = cv2.VideoCapture(input_file)
    if not cap.isOpened():
        raise RuntimeError(f"无法打开视频文件: {input_file}")
    fps = cap.get(cv2.CAP_PROP_FPS)
    if not fps or fps < 1:
        fps = 25
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration_sec = total_frames / fps if fps > 0 else 0
    segments = []
    video_basename = os.path.splitext(os.path.basename(input_file))[0]
    idx = 0
    segment_start_sec = 0
    while segment_start_sec < duration_sec:
        cap.set(cv2.CAP_PROP_POS_FRAMES, int(segment_start_sec * fps))
        ret, frame = cap.read()
        img_path = os.path.join(PUBLIC_DIR, f"{video_basename}_keyframe_{idx+1}.jpg")
        if ret:
            cv2.imwrite(img_path, frame)
        else:
            img_path = None
        audio_path = os.path.join(PUBLIC_DIR, f"{video_basename}_seg_{idx+1}.wav")
        seg_dur = min(force_segment_sec, duration_sec - segment_start_sec)
        cmd = [
            "ffmpeg", "-y", "-i", input_file,
            "-ss", str(segment_start_sec),
            "-t", str(seg_dur),
            "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", audio_path
        ]
        try:
            subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        except Exception:
            audio_path = None
        if audio_path and not os.path.exists(audio_path):
            audio_path = None
        segments.append({
            "index": idx,
            "text": None,
            "audio": audio_path,
            "image": [img_path] if img_path else []
        })
        idx += 1
        segment_start_sec += force_segment_sec
    cap.release()
    if not segments:
        cap = cv2.VideoCapture(input_file)
        ret, frame = cap.read()
        img_path = os.path.join(PUBLIC_DIR, f"{video_basename}_keyframe_1.jpg")
        if ret:
            cv2.imwrite(img_path, frame)
        else:
            img_path = None
        audio_path = os.path.join(PUBLIC_DIR, f"{video_basename}_seg_1.wav")
        cmd = [
            "ffmpeg", "-y", "-i", input_file,
            "-ss", "0", "-t", str(force_segment_sec),
            "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", audio_path
        ]
        try:
            subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        except Exception:
            audio_path = None
        if audio_path and not os.path.exists(audio_path):
            audio_path = None
        segments.append({
            "index": 0,
            "text": None,
            "audio": audio_path,
            "image": [img_path] if img_path else []
        })
        cap.release()
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(segments, f, ensure_ascii=False, indent=2)
    print(f"[VIDEO] 分段完成，共{len(segments)}段，输出: {output_path}")
# 只以画面变化为唯一分段依据，采用 pySceneDetect 新API，兼容 Python 3.10+，无需手动调参。

def audio_asr_and_refine_task(input_json, output_json=None, config_path=None):
    print(f"[ASR+LLM] 开始音频ASR识别与润色: {input_json}")
    """
    对音频分段json做ASR+LLM润色，输出新json。audio字段只保留ASR文本，不能识别则置空。text字段保持原样不动。
    """
    import os, json, yaml
    if not os.path.exists(input_json):
        raise FileNotFoundError(f"{input_json} not found")
    with open(input_json, 'r', encoding='utf-8') as f:
        segments = json.load(f)
    if not config_path:
        config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    try:
        import whisper
        model = whisper.load_model(config.get('asr_model', 'base'))
    except Exception:
        model = None
    refined_segments = []
    for seg in segments:
        audio_path = seg.get('audio')
        # 只处理audio字段，text字段不动
        if not audio_path or not os.path.exists(audio_path) or not model:
            seg['audio'] = ''
            refined_segments.append(seg)
            continue
        try:
            result = model.transcribe(audio_path, language=config.get('asr_language', 'zh'))
            asr_text = result['text'].strip()
        except Exception:
            asr_text = ''
        seg['audio'] = asr_text if asr_text else ''
        refined_segments.append(seg)
    if not output_json:
        output_json = os.path.splitext(input_json)[0] + '_refined.json'
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(refined_segments, f, ensure_ascii=False, indent=2)
    print(f"[ASR+LLM] 识别与润色完成，输出: {output_json}")

def video_multimodal_understanding_task(input_json, output_json=None, config_path=None):
    """
    对视频分段json做多模态理解，输出新json。无Airflow依赖，参数直接传递。
    """
    import os, json, yaml
    if not os.path.exists(input_json):
        raise FileNotFoundError(f"{input_json} not found")
    with open(input_json, 'r', encoding='utf-8') as f:
        segments = json.load(f)
    if not config_path:
        config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    multimodal_api = config.get('multimodal_api', 'http://localhost:8000/multimodal_infer')
    multimodal_model = config.get('multimodal_model', 'qwen-vl')
    try:
        import whisper
        asr_model = whisper.load_model(config.get('asr_model', 'base'))
    except Exception:
        asr_model = None
    refined_segments = []
    for seg in segments:
        img_path = seg['image'][0] if seg['image'] else None
        audio_path = seg['audio']
        asr_text = None
        if audio_path and asr_model and os.path.exists(audio_path):
            try:
                result = asr_model.transcribe(audio_path, language=config.get('asr_language', 'zh'))
                asr_text = result['text'].strip()
            except Exception:
                asr_text = None
        multimodal_result = None
        if img_path and os.path.exists(img_path):
            try:
                multimodal_result = {
                    'summary': f"[DEMO] 图像+音频内容摘要 (index={seg['index']})",
                    'tags': ["demo_tag1", "demo_tag2"],
                    'asr_text': asr_text,
                    'model': multimodal_model
                }
            except Exception as e:
                multimodal_result = {'error': str(e)}
        seg['multimodal_result'] = multimodal_result
        seg['text'] = asr_text
        refined_segments.append(seg)
    if not output_json:
        output_json = os.path.splitext(input_json)[0] + '_multimodal.json'
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(refined_segments, f, ensure_ascii=False, indent=2)
    print(f"视频多模态理解完成，输出: {output_json}")

def asr_whisper_for_segments(segments, config=None):
    """
    对分段列表中的 audio 字段（如为音频文件路径）用 whisper 识别为文本，直接覆盖 audio 字段（即 audio 字段只保留文本，不能识别则置空）。
    """
    import os
    try:
        import whisper
    except ImportError:
        print("[ASR] 未安装 whisper，跳过音频识别。请 pip install -U openai-whisper")
        return segments
    model_name = (config or {}).get('asr_model', 'base')
    model = whisper.load_model(model_name)
    for seg in segments:
        audio_val = seg.get('audio')
        if isinstance(audio_val, str) and audio_val.lower().endswith((".wav", ".mp3", ".m4a", ".flac", ".aac", ".ogg")) and os.path.exists(audio_val):
            try:
                result = model.transcribe(audio_val, language=(config or {}).get('asr_language', 'zh'))
                seg['audio'] = result['text'].strip() if result['text'].strip() else ''
            except Exception as e:
                print(f"[ASR ERROR] {audio_val}: {e}")
                seg['audio'] = ''
        elif isinstance(audio_val, str) and audio_val.strip() and not audio_val.lower().endswith((".wav", ".mp3", ".m4a", ".flac", ".aac", ".ogg")):
            # 已经是文本，保留
            continue
        else:
            seg['audio'] = ''
    return segments

def refine_text_for_segments(json_path, output_path=None, config=None):
    print(f"[LLM] 开始文本润色: {json_path}")
    """
    对分段json文件中的每个segment的text字段进行自然语言优化（如错别字纠正、润色、可读性提升），输出新json。
    自动读取 config.yaml 里的 llm 配置，调用大模型API。
    """
    import os, json
    import requests
    import yaml
    # 读取配置
    if config is None:
        config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
        with open(config_path, 'r') as f:
            config = yaml.safe_load(f)
    llm_cfg = config.get('llm', {})
    llm_api = llm_cfg.get('api') or llm_cfg.get('llm_api')
    llm_model = llm_cfg.get('model') or llm_cfg.get('llm_model')
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"{json_path} not found")
    with open(json_path, 'r', encoding='utf-8') as f:
        segments = json.load(f)
    def refine_text_with_llm(text):
        if not llm_api:
            return text
        try:
            resp = requests.post(llm_api, json={
                'model': llm_model,
                'prompt': text
            }, timeout=30)
            if resp.status_code == 200:
                data = resp.json()
                # 假设返回格式为 {'result': '优化后文本'}
                return data.get('result', text)
            else:
                return text
        except Exception as e:
            print(f"[LLM ERROR] {e}")
            return text
    refined_segments = []
    for seg in segments:
        text = seg.get('text')
        if text and isinstance(text, str) and text.strip():
            try:
                refined_text = refine_text_with_llm(text)
            except Exception:
                refined_text = text
            seg['text'] = refined_text
        refined_segments.append(seg)
    if not output_path:
        base, ext = os.path.splitext(json_path)
        output_path = base + '_refined' + ext
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(refined_segments, f, ensure_ascii=False, indent=2)
    print(f"[LLM] 文本润色完成，输出: {output_path}")

def ocr_image_for_segments(json_path, output_path=None, config=None):
    print(f"[OCR] 开始图片OCR识别: {json_path}")
    """
    对分段json文件中的每个segment的image字段（图片路径列表）做OCR，识别结果只覆盖image字段为文字（str），其它字段（如text、audio等）必须保持原样不动。
    不增加新字段，功能单一、节点解耦。
    """
    import os, json
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        print("[OCR] 未安装 pytesseract 或 pillow，跳过图片转文字。请 pip install pytesseract pillow")
        return
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"{json_path} not found")
    with open(json_path, 'r', encoding='utf-8') as f:
        segments = json.load(f)
    def ocr_image(img_path):
        if not os.path.exists(img_path):
            return ''
        try:
            # 兼容中英文，优先用chi_sim+eng
            try:
                return pytesseract.image_to_string(Image.open(img_path), lang="chi_sim+eng").strip()
            except Exception:
                return pytesseract.image_to_string(Image.open(img_path), lang="eng").strip()
        except Exception as e:
            print(f"[OCR ERROR] {img_path}: {e}")
            return ''
    refined_segments = []
    for seg in segments:
        # 只处理 image 字段，其他字段完全不动
        img_list = seg.get('image') or []
        ocr_texts = []
        for img_path in img_list:
            ocr_result = ocr_image(img_path)
            if ocr_result:
                ocr_texts.append(ocr_result)
        # 只覆盖 image 字段为OCR文本，其他字段保持原样
        new_seg = dict(seg)
        new_seg['image'] = "\n".join(ocr_texts) if ocr_texts else ''
        refined_segments.append(new_seg)
    if not output_path:
        base, ext = os.path.splitext(json_path)
        output_path = base + '_ocr' + ext
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(refined_segments, f, ensure_ascii=False, indent=2)
    print(f"[OCR] 图片OCR识别完成，输出: {output_path}")

def fuse_segments_to_text(json_path, output_path=None, use_llm=True, config=None):
    print(f"[FUSE] 开始融合所有分段文本: {json_path}")
    """
    融合分段json的text/audio/image三个字段为一个字符串，可选用大模型进一步摘要/润色。
    - 只输出有意义的内容，彻底去除音频、图片等文件路径。
    """
    import os, json
    import yaml
    import requests
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"{json_path} not found")
    with open(json_path, 'r', encoding='utf-8') as f:
        segments = json.load(f)
    fused_list = []
    for seg in segments:
        # 只保留有意义的内容，不输出任何文件路径
        parts = []
        text = seg.get('text')
        if text and isinstance(text, str) and text.strip():
            parts.append(text.strip())
        else:
            # audio字段如为转写文本直接用，否则忽略
            audio_val = seg.get('audio')
            if audio_val and isinstance(audio_val, str) and not audio_val.lower().endswith(('.wav','.mp3','.m4a','.flac','.aac','.ogg')):
                parts.append(audio_val.strip())
            # image字段如为文字直接用，否则忽略
            image_val = seg.get('image')
            if image_val:
                if isinstance(image_val, str) and image_val.strip() and not image_val.lower().endswith(('.jpg','.jpeg','.png','.bmp','.gif','.tiff')):
                    parts.append(image_val.strip())
                elif isinstance(image_val, list):
                    for img in image_val:
                        if isinstance(img, str) and img.strip() and not img.lower().endswith(('.jpg','.jpeg','.png','.bmp','.gif','.tiff')):
                            parts.append(img.strip())
        if parts:
            fused_list.append("\n".join(parts))
    fused_text = "\n---\n".join(fused_list)
    # 是否用大模型进一步处理
    if use_llm:
        if config is None:
            config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
            with open(config_path, 'r') as f:
                config = yaml.safe_load(f)
        llm_cfg = config.get('llm', {})
        llm_api = llm_cfg.get('api') or llm_cfg.get('llm_api')
        llm_model = llm_cfg.get('model') or llm_cfg.get('llm_model')
        if llm_api:
            try:
                resp = requests.post(llm_api, json={
                    'model': llm_model,
                    'prompt': fused_text
                }, timeout=60)
                if resp.status_code == 200:
                    data = resp.json()
                    fused_text = data.get('result', fused_text)
            except Exception as e:
                print(f"[LLM FUSE ERROR] {e}")
    if not output_path:
        base, ext = os.path.splitext(json_path)
        output_path = base + '_fused.txt'
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(fused_text)
    print(f"[FUSE] 融合输出完成，输出: {output_path}")

def run_pipeline(id: str):
    """
    多模态主流程入口（ID驱动）。严格顺序：分段->ASR->大模型润色->OCR->融合输出。
    只通过config.yaml的public_dir存取所有中间产物。
    上传即生成唯一ID，所有分段/片段/中间产物都以ID为前缀命名，所有节点只传递ID参数，所有内容都覆盖写回ID.json。
    """
    import os, yaml, glob
    import time
    import json
    # 读取全局配置
    config_path = os.path.join(os.path.dirname(__file__), 'config/config.yaml')
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    public_dir = config.get('public_dir', './public')
    os.makedirs(public_dir, exist_ok=True)
    # 1. 查找原始文件（ID.扩展名）
    file_list = list(glob.glob(os.path.join(public_dir, f"{id}.*")))
    if not file_list:
        raise FileNotFoundError(f"未找到ID对应的原始文件: {id}")
    input_file = file_list[0]
    # 2. 自动识别类型并分段，分段json为ID.json
    seg_func_map = {
        'segment_text': segment_text_task,
        'segment_pdf': segment_pdf_task,
        'segment_word': segment_word_task,
        'segment_ppt': segment_ppt_task,
        'segment_excel': segment_excel_task,
        'segment_audio': segment_audio_task,
        'segment_video': segment_video_task,
    }
    seg_node = classify_file_task(input_file)
    print(f"[PIPELINE] 文件类型识别为: {seg_node}")
    seg_json = os.path.join(public_dir, f"{id}.json")
    seg_func = seg_func_map.get(seg_node)
    if not seg_func:
        raise ValueError(f"不支持的文件类型: {seg_node}")
    # 分段/片段文件全部以ID_001.xxx命名，需在各分段节点实现
    seg_func(input_file, seg_json)
    # 3. ASR音频识别（严格在OCR前），覆盖写回ID.json
    print(f"[PIPELINE] 开始ASR音频识别: {seg_json}")
    audio_asr_and_refine_task(seg_json, seg_json, config_path)
    # 4. 可选：大模型润色text字段，覆盖写回ID.json
    print(f"[PIPELINE] 开始大模型润色: {seg_json}")
    refine_text_for_segments(seg_json, seg_json, config)
    # 5. 图片OCR识别（必须在ASR后），覆盖写回ID.json
    print(f"[PIPELINE] 开始图片OCR识别: {seg_json}")
    ocr_image_for_segments(seg_json, seg_json, config)
    # 6. 融合输出，输出为ID_fused.txt
    fused_txt = os.path.join(public_dir, f"{id}_fused.txt")
    print(f"[PIPELINE] 开始融合输出: {seg_json} -> {fused_txt}")
    fuse_segments_to_text(seg_json, fused_txt, use_llm=True, config=config)
    print(f"[PIPELINE] 全流程完成，最终输出: {fused_txt}")
